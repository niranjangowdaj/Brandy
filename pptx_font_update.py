import os
from pptx import Presentation
from pptx.util import Pt, Inches


FONT_NAME = "Arial" 
FONT_SIZE = 24       
FOOTER_TEXT = "© 2023 SAP SE or an SAP affiliate company. All rights reserved | PUBLIC"
FOOTER_FONT = "Times New Roman"
FOOTER_SIZE = 10

def add_footer_to_slide(slide, text):
    left = Inches(0.2) 
    width = Inches(8)  
    height = Inches(0.4)
    slide_height = slide.part.slide_layout.slide_height if hasattr(slide.part.slide_layout, 'slide_height') else Inches(7.5)
    top = Inches(7.0) 
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FOOTER_FONT
    run.font.size = Pt(FOOTER_SIZE)

pptx_file = None
for file in os.listdir('.'):
    if file.lower().endswith('.pptx'):
        pptx_file = file
        break

if not pptx_file:
    print("No .pptx file found in the current directory.")
    exit(1)

prs = Presentation(pptx_file)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(FONT_SIZE)
    add_footer_to_slide(slide, FOOTER_TEXT)

output_file = f"modified_{pptx_file}"
prs.save(output_file)
print(f"Saved modified presentation as {output_file}") 