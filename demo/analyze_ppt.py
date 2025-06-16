#!/usr/bin/env python3
"""
PowerPoint Element Analyzer
Prints all elements from each slide in a PowerPoint file to the terminal.

Usage: python analyze_ppt.py <path_to_pptx_file>
"""

import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_shape(shape, shape_idx):
    details = []
    
    
    details.append(f"  Element {shape_idx}:")
    details.append(f"    Type: {type(shape).__name__}")
    details.append(f"    Shape ID: {shape.shape_id}")
    
   
    try:
        details.append(f"    Position: ({shape.left.inches:.2f}, {shape.top.inches:.2f}) inches")
        details.append(f"    Size: {shape.width.inches:.2f} x {shape.height.inches:.2f} inches")
    except AttributeError:
        details.append("    Position/Size: Not available")
    

    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        text_content = shape.text.strip()
        if text_content:
            details.append(f"    Text: '{text_content}'")
            
     
            try:
                fonts_found = set()
                font_sizes = set()
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():  
                            font_name = run.font.name if run.font.name else "Default"
                            fonts_found.add(font_name)
                            
                            if run.font.size:
                                font_sizes.add(str(run.font.size.pt))
                            else:
                                font_sizes.add("Default")
                
                if fonts_found:
                    details.append(f"    Fonts: {', '.join(sorted(fonts_found))}")
                if font_sizes:
                    details.append(f"    Font Sizes: {', '.join(sorted(font_sizes))}pt")
                    
            except Exception as e:
                try:
                    first_paragraph = shape.text_frame.paragraphs[0]
                    if first_paragraph.runs:
                        first_run = first_paragraph.runs[0]
                        font_name = first_run.font.name or "Default"
                        font_size = first_run.font.size.pt if first_run.font.size else "Default"
                        details.append(f"    Font: {font_name}, Size: {font_size}pt")
                except:
                    details.append("    Font: Unable to detect")
        else:
            details.append("    Text: [Empty text frame]")
    

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            details.append("    Content: Image/Picture")
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                details.append(f"    Content: Table ({shape.table.rows} rows x {shape.table.columns} cols)")
            except:
                details.append("    Content: Table")
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            details.append("    Content: Chart/Graphic")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            details.append("    Content: Group of shapes")
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            details.append("    Content: AutoShape")
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            details.append("    Content: Text Box")
        else:
            details.append(f"    Shape Type: {shape.shape_type}")
    except AttributeError:
        details.append("    Shape Type: Unknown")
    
    return "\n".join(details)

def analyze_pptx(file_path):
    """Analyze PowerPoint file and print all elements"""
    try:
        prs = Presentation(file_path)
        print(f"\n=== PowerPoint Analysis: {file_path} ===")
        print(f"Total slides: {len(prs.slides)}\n")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f"SLIDE {slide_idx}:")
            print(f"  Layout: {slide.slide_layout.name}")
            print(f"  Total elements: {len(slide.shapes)}")
            
            if len(slide.shapes) == 0:
                print("  [No elements found]")
            else:
                for shape_idx, shape in enumerate(slide.shapes, 1):
                    print(analyze_shape(shape, shape_idx))
            
            # Check for speaker notes
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                print(f"  Speaker Notes: '{notes_text}'")
            
            print("-" * 60)
        
        print(f"\nAnalysis complete! Processed {len(prs.slides)} slides.")
        
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
        return False
    except Exception as e:
        print(f"Error analyzing PowerPoint file: {e}")
        return False
    
    return True

def main():
    if len(sys.argv) != 2:
        print("Usage: python analyze_ppt.py <path_to_pptx_file>")
        print("Example: python analyze_ppt.py presentation.pptx")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    success = analyze_pptx(pptx_file)
    
    if not success:
        sys.exit(1)

if __name__ == "__main__":
    main() 