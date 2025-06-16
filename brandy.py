import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from docx import Document
import io
import os
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
import numpy as np
import json
import re
import pandas as pd
from sklearn.metrics.pairwise import cosine_similarity

st.set_page_config(page_title="Brandy", layout="wide", initial_sidebar_state="expanded")

# Header with title and clear chat button
col1, col2 = st.columns([4, 1])
with col1:
    st.title("Brandy - Your Brand Assistant")
with col2:
    st.write("")  # Add some spacing
    if st.button("🗑️ Clear Chat", help="Clear chat history"):
        st.session_state.chat_history = []
        st.rerun()

FOOTER_FONT = "SAP 72"
FOOTER_SIZE = 8

BRAND_GUIDELINES_PATH = "Project Brandy - Brand Guidelines for PPTs.docx"

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "docx_text" not in st.session_state:
    st.session_state.docx_text = None
if "doc_chunks" not in st.session_state:
    st.session_state.doc_chunks = None
if "chunk_embeddings" not in st.session_state:
    st.session_state.chunk_embeddings = None
if "pptx_issues" not in st.session_state:
    st.session_state.pptx_issues = None
if "pptx_modified" not in st.session_state:
    st.session_state.pptx_modified = None
if "sentence_model" not in st.session_state:
    st.session_state.sentence_model = None
if "gemini_model" not in st.session_state:
    st.session_state.gemini_model = None
if "links_df" not in st.session_state:
    st.session_state.links_df = None
if "links_embeddings" not in st.session_state:
    st.session_state.links_embeddings = None

# Gemini API key
gemini_api_key = ""

def load_embeddings_and_chunks(prefix):
    data = np.load(f"{prefix}_embeddings.npz")
    with open(f"{prefix}_chunks.json") as f:
        chunks = json.load(f)
    return data['embeddings'], chunks

def find_relevant_links(answer_text, top_k=3):
    if (st.session_state.links_df is None or 
        st.session_state.links_embeddings is None or 
        st.session_state.sentence_model is None):
        return []
    
    answer_embedding = st.session_state.sentence_model.encode([answer_text])
    
    similarities = cosine_similarity(answer_embedding, st.session_state.links_embeddings)[0]
    
    top_indices = similarities.argsort()[-top_k:][::-1]
    
    relevant_links = []
    for idx in top_indices:
        if similarities[idx] > 0.3: 
            relevant_links.append({
                'name': st.session_state.links_df.iloc[idx]['Name'],
                'link': st.session_state.links_df.iloc[idx]['Link'],
                'similarity': similarities[idx]
            })
    
    return relevant_links

def display_relevant_links(links):
    if not links:
        return
    
    cols = st.columns(min(len(links), 3))
    
    for i, link_info in enumerate(links):
        with cols[i % 3]:
            st.markdown(f"""
                <a href="{link_info['link']}" target="_blank" style="text-decoration: none;">
                    <div style="
                        background-color: transparent;
                        border: 1px solid white;
                        border-radius: 6px;
                        padding: 4px 8px;
                        margin: 2px 0;
                        text-align: center;
                        color: white;
                        font-size: 11px;
                        transition: all 0.3s ease;
                        cursor: pointer;
                        opacity: 0.7;
                        width: fit-content;
                    " onmouseover="this.style.boxShadow='0 0 8px rgba(255,255,255,0.3)'; this.style.opacity='1';" 
                       onmouseout="this.style.boxShadow='none'; this.style.opacity='0.7';">
                        {link_info['name']}
                    </div>
                </a>
            """, unsafe_allow_html=True)

if os.path.exists(BRAND_GUIDELINES_PATH):
    try:
        doc = Document(BRAND_GUIDELINES_PATH)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        st.session_state.docx_text = "\n".join(full_text)
        
        prefix = "mydoc"
        embeddings, chunks = load_embeddings_and_chunks(prefix)
        st.session_state.chunk_embeddings = embeddings
        st.session_state.doc_chunks = chunks
        
        if st.session_state.sentence_model is None:
            with st.spinner("Loading..."):
                st.session_state.sentence_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        if st.session_state.gemini_model is None:
            genai.configure(api_key=gemini_api_key)
            st.session_state.gemini_model = genai.GenerativeModel('gemini-1.5-flash-002')
        
        if st.session_state.links_df is None and os.path.exists("links.csv"):
            st.session_state.links_df = pd.read_csv("links.csv")
            link_names = st.session_state.links_df['Name'].tolist()
            st.session_state.links_embeddings = st.session_state.sentence_model.encode(link_names)
            
    except Exception as e:
        st.sidebar.error(f"Error loading brand guidelines or embeddings: {str(e)}")
else:
    st.sidebar.error("Brand guidelines document not found!")

st.sidebar.header("Upload Files")
uploaded_file = st.sidebar.file_uploader("Upload file for compliance check", type=["pptx", "docx", "pdf"])

def handle_pdf_compliance(file):
    st.sidebar.info("PDF compliance check coming soon!")
    return None, None

def handle_docx_compliance(file):
    st.sidebar.info("DOCX compliance check coming soon!")
    return None, None

def add_red_border(shape):
    try:
        line = shape.line
        line.color.rgb = RGBColor(255, 0, 0)
        line.width = Pt(3)
    except Exception:
        pass

def add_green_border(shape):
    try:
        line = shape.line
        line.color.rgb = RGBColor(0, 255, 0) 
        line.width = Pt(3)
    except Exception:
        pass

def add_footer_to_slide(slide, text):
    left = Inches(0.2) 
    width = Inches(8)  
    height = Inches(0.4)
    slide_height = slide.part.slide_layout.slide_height if hasattr(slide.part.slide_layout, 'slide_height') else Inches(7.5)
    top = Inches(7.0) 
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FOOTER_FONT
    run.font.size = Pt(FOOTER_SIZE)

def add_summary_slide(prs, issues):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "PPTX Compliance Issues"
    body = "\n".join(issues) if issues else "No issues found."
    if len(slide.shapes) > 1:
        slide.shapes[1].text = body
    else:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = body
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.insert(0, slides[-1])
    xml_slides.remove(slides[-1])

def is_sentence_case(text):
    if not text or not text.strip():
        return False
    stripped = text.lstrip()
    match = re.search(r'[A-Za-z]', stripped)
    if not match:
        return False
    first_alpha_idx = match.start()
    if not stripped[first_alpha_idx].isupper():
        return False
    words = stripped.split()
    if len(words) > 1:
        for word in words[1:]:
            if len(word) > 1 and word.isupper():
                continue
            if word and word[0].isupper():
                return False
    return True

def check_element_compliance(element_info, gemini_model):
    prompt = f"""
    Given these brand guidelines for PowerPoint presentations:
    {st.session_state.docx_text}
    
    Check if this element complies with the guidelines:
    {element_info}
    
    Respond with either:
    - "COMPLIANT" if the element follows all relevant guidelines
    - "NON-COMPLIANT: [specific reason]" if it violates any guidelines
    """
    
    response = gemini_model.generate_content(prompt)
    answer = response.text.strip()
    
    is_compliant = answer.startswith("COMPLIANT")
    return is_compliant, answer

def pptx_compliance_check_with_rules(pptx_file, rules, add_copyright, copyright_type, implement_actions=False):
    prs = Presentation(pptx_file)
    issues = []
    slide_comments = {}
    
    if add_copyright:
        copyright_text = "© SAP SE or an SAP affiliate company. All rights reserved. Internal Use Only." if copyright_type == "Internal" else "© SAP SE or an SAP affiliate company. All rights reserved. Public Use."
        for slide in prs.slides:
            add_footer_to_slide(slide, copyright_text)
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_issue_comments = []
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if not shape.has_text_frame:
                continue
            
            # Collect all element information
            element_info = {
                "slide_number": slide_idx,
                "element_number": shape_idx,
                "text": shape.text,
                "text_case": {
                    "is_uppercase": shape.text.isupper(),
                    "is_lowercase": shape.text.islower(),
                    "is_title_case": shape.text.istitle(),
                    "is_sentence_case": is_sentence_case(shape.text)
                },
                "font_details": []
            }
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "©" in run.text:
                        break
                    
                    current_size = None
                    try:
                        if run.font.size:
                            current_size = run.font.size.pt
                    except Exception:
                        pass
                    
                    if implement_actions:
                        run.font.name = "SAP 72"
                        
                        if current_size is not None and current_size < 11:
                            run.font.size = Pt(11)
                            add_green_border(shape)
                            slide_issue_comments.append(f"Element {shape_idx}: Font size increased to 11pt")
                    
                    font_info = {
                        "font_name": run.font.name,
                        "font_size": current_size,
                        "text": run.text,
                        "text_case": {
                            "is_uppercase": run.text.isupper(),
                            "is_lowercase": run.text.islower(),
                            "is_title_case": run.text.istitle(),
                            "is_sentence_case": is_sentence_case(run.text)
                        }
                    }
                    element_info["font_details"].append(font_info)
            
            if("©" not in element_info["text"] and element_info["text"] != ""):
                is_compliant, compliance_message = check_element_compliance(str(element_info), st.session_state.gemini_model)
                
                if not is_compliant:
                    issues.append(f"Slide {slide_idx}, Element {shape_idx}: {compliance_message}")
                    slide_issue_comments.append(f"Element {shape_idx}: {compliance_message}")
                    add_red_border(shape)
            else:
                is_compliant = True
                compliance_message = "Compliant"
        
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        if slide_issue_comments:
            notes_text_frame.text = f"Slide {slide_idx} compliance issues:\n" + "\n".join(slide_issue_comments)
        else:
            notes_text_frame.text = f"Slide {slide_idx}: All elements compliant."
    
    add_summary_slide(prs, issues)
    
    temp_file = io.BytesIO()
    prs.save(temp_file)
    temp_file.seek(0)
    
    prs = Presentation(temp_file)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return issues, output

if uploaded_file:
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    if file_type == 'pptx':
        st.sidebar.subheader("PPTX Options")
        add_copyright = st.sidebar.checkbox("Add Copyright Footer", value=True)
        
        if add_copyright:
            copyright_type = st.sidebar.selectbox(
                "Copyright Type",
                ["Internal", "Public"],
                help="Select whether this is for internal or public use"
            )
        
        implement_actions = st.sidebar.checkbox(
            "Implement Actions",
            help="Automatically fix font and size issues (SAP 72 font and minimum 11pt size)"
        )
        
        if st.sidebar.button("Run Compliance Check"):
            with st.spinner("Checking PPTX compliance..."):
                issues, pptx_bytes = pptx_compliance_check_with_rules(
                    uploaded_file, 
                    "",  
                    add_copyright,
                    copyright_type if add_copyright else None,
                    implement_actions
                )
                st.session_state.pptx_issues = issues
                st.session_state.pptx_modified = pptx_bytes
            st.sidebar.success("Compliance check complete! Download the modified PPTX below.")
            
            if st.session_state.pptx_modified:
                st.sidebar.download_button(
                    label="Download Modified PPTX",
                    data=st.session_state.pptx_modified,
                    file_name="pptx_compliance_checked.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                # it doesn't look good
                # if st.session_state.pptx_issues:
                #     st.sidebar.info("\n".join(st.session_state.pptx_issues[:10]) + ("\n..." if len(st.session_state.pptx_issues) > 10 else ""))
    
    elif file_type == 'pdf':
        handle_pdf_compliance(uploaded_file)
    
    elif file_type == 'docx':
        handle_docx_compliance(uploaded_file)

if user_input := st.chat_input("Ask a question about brand guidelines..."):
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    
    if (st.session_state.doc_chunks and 
        st.session_state.chunk_embeddings is not None and 
        st.session_state.sentence_model is not None and
        st.session_state.gemini_model is not None):
        
        relevant_chunks = []
        q_emb = st.session_state.sentence_model.encode([user_input])[0]
        similarities = np.dot(st.session_state.chunk_embeddings, q_emb)
        top_indices = similarities.argsort()[-3:][::-1]
        for i in top_indices:
            relevant_chunks.append(st.session_state.doc_chunks[i])
        
        context = "\n\n".join(relevant_chunks)
        prompt = f"Answer the question based on the following context from SAP brand guidelines:\n\n{context}\n\nQuestion: {user_input}"
        response = st.session_state.gemini_model.generate_content(prompt)
        answer = response.text.strip()
        
        st.session_state.chat_history.append({"role": "assistant", "content": answer})
        
    else:
        error_msg = "Please ensure brand guidelines and AI model are loaded to enable Q&A functionality."
        st.session_state.chat_history.append({"role": "assistant", "content": error_msg})

for i, msg in enumerate(st.session_state.chat_history):
    with st.chat_message(msg["role"]):
        st.write(msg["content"])
        
        if msg["role"] == "assistant":
            relevant_links = find_relevant_links(msg["content"])
            if relevant_links:
                display_relevant_links(relevant_links) 