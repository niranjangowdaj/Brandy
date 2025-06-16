import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from docx import Document
import io
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
import numpy as np
import json

st.set_page_config(page_title="SAP Compliance & Q&A Bot", layout="wide")
st.title("SAP Compliance & Q&A Bot (MVP)")


if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "docx_text" not in st.session_state:
    st.session_state.docx_text = None
if "doc_chunks" not in st.session_state:
    st.session_state.doc_chunks = None
if "chunk_embeddings" not in st.session_state:
    st.session_state.chunk_embeddings = None
if "pptx_issues" not in st.session_state:
    st.session_state.pptx_issues = None
if "pptx_modified" not in st.session_state:
    st.session_state.pptx_modified = None


gemini_api_key = ""

st.sidebar.header("Upload Reference Files")
docx_file = st.sidebar.file_uploader("Upload Word Document for Q&A & Compliance (for compliance only, not for embedding)", type=["docx"])
pptx_file = st.sidebar.file_uploader("Upload PPTX for Compliance", type=["pptx"])


def load_embeddings_and_chunks(prefix):
    data = np.load(f"{prefix}_embeddings.npz")
    with open(f"{prefix}_chunks.json") as f:
        chunks = json.load(f)
    return data['embeddings'], chunks

prefix = "mydoc"
if prefix:
    try:
        embeddings, chunks = load_embeddings_and_chunks(prefix)
        st.session_state.chunk_embeddings = embeddings
        st.session_state.doc_chunks = chunks
        st.sidebar.success(f"Loaded {len(chunks)} chunks and embeddings for prefix '{prefix}'.")
    except Exception as e:
        st.sidebar.error(f"Failed to load embeddings/chunks: {e}")

# Extract text from docx (for compliance rules)
def extract_docx_text(docx_file):
    doc = Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

# Extract compliance rules from docx using Gemini
def extract_compliance_rules(docx_text, gemini_api_key):
    genai.configure(api_key=gemini_api_key)
    model = genai.GenerativeModel('gemini-1.5-flash-002')
    prompt = (
        "Extract a concise, clear list of compliance rules for PowerPoint presentations from the following document. "
        "Format as numbered rules, each on a new line.\n\nDocument:\n" + docx_text
    )
    response = model.generate_content(prompt)
    rules = response.text.strip()
    return rules

# Add red border to a shape
def add_red_border(shape):
    try:
        line = shape.line
        line.color.rgb = RGBColor(255, 0, 0)
        line.width = Pt(3)
    except Exception:
        pass

# Add summary slide
def add_summary_slide(prs, issues):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "PPTX Compliance Issues"
    body = "\n".join(issues) if issues else "No issues found."
    if len(slide.shapes) > 1:
        slide.shapes[1].text = body
    else:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = body
    # Move summary slide to first position
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.insert(0, slides[-1])
    xml_slides.remove(slides[-1])

# Compliance check using rules from docx
def pptx_compliance_check_with_rules(pptx_file, rules, gemini_api_key):
    prs = Presentation(pptx_file)
    issues = []
    slide_comments = {}
    genai.configure(api_key=gemini_api_key)
    model = genai.GenerativeModel('gemini-1.5-flash-002')
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_issue_comments = []
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if not shape.has_text_frame:
                continue
            text = shape.text
            prompt = (
                f"Given these compliance rules for PowerPoint presentations:\n{rules}\n\n"
                f"Is the following slide element compliant?\nElement text: {text}\n"
                f"If not, explain why in one sentence. If compliant, reply 'Compliant'."
            )
            response = model.generate_content(prompt)
            answer = response.text.strip()
            if answer.lower() != "compliant":
                issues.append(f"Slide {slide_idx}, Element {shape_idx}: {answer}")
                slide_issue_comments.append(f"Element {shape_idx}: {answer}")
                add_red_border(shape)
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        if slide_issue_comments:
            notes_text_frame.text = f"Slide {slide_idx} compliance issues:\n" + "\n".join(slide_issue_comments)
        else:
            notes_text_frame.text = f"Slide {slide_idx}: All elements compliant."
    add_summary_slide(prs, issues)
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return issues, output

# Handle docx upload (for compliance rules only)
if docx_file:
    st.session_state.docx_text = extract_docx_text(docx_file)
    st.sidebar.success("Word document loaded for compliance rule extraction.")

# Sidebar: PPTX Compliance Checker
st.sidebar.header("PPTX Compliance Checker")
if st.sidebar.button("Run Compliance Check on PPTX") and pptx_file and st.session_state.docx_text and gemini_api_key:
    with st.spinner("Extracting compliance rules and checking PPTX..."):
        rules = extract_compliance_rules(st.session_state.docx_text, gemini_api_key)
        issues, pptx_bytes = pptx_compliance_check_with_rules(pptx_file, rules, gemini_api_key)
        st.session_state.pptx_issues = issues
        st.session_state.pptx_modified = pptx_bytes
    st.sidebar.success("Compliance check complete! Download the modified PPTX below.")

if st.session_state.pptx_modified:
    st.sidebar.download_button(
        label="Download Modified PPTX",
        data=st.session_state.pptx_modified,
        file_name="pptx_compliance_checked.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    if st.session_state.pptx_issues:
        st.sidebar.info("\n".join(st.session_state.pptx_issues[:10]) + ("\n..." if len(st.session_state.pptx_issues) > 10 else ""))

# Chat interface
st.subheader("Chat with your reference document or get compliance results:")
chat_container = st.container()

for msg in st.session_state.chat_history:
    if msg["role"] == "user":
        chat_container.markdown(f"**You:** {msg['content']}")
    else:
        chat_container.markdown(f"**Bot:** {msg['content']}")

user_input = st.chat_input("Ask a question about the loaded document...")

if user_input and st.session_state.doc_chunks and st.session_state.chunk_embeddings is not None and gemini_api_key:
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    model = SentenceTransformer('all-MiniLM-L6-v2')
    relevant_chunks = []
    # Use loaded embeddings and chunks
    q_emb = model.encode([user_input])[0]
    similarities = np.dot(st.session_state.chunk_embeddings, q_emb)
    top_indices = similarities.argsort()[-3:][::-1]
    for i in top_indices:
        relevant_chunks.append(st.session_state.doc_chunks[i])
    context = "\n\n".join(relevant_chunks)
    genai.configure(api_key=gemini_api_key)
    gemini_model = genai.GenerativeModel('gemini-1.5-flash-002')
    prompt = f"Answer the question based on the following context:\n\n{context}\n\nQuestion: {user_input}"
    response = gemini_model.generate_content(prompt)
    answer = response.text.strip()
    st.session_state.chat_history.append({"role": "assistant", "content": answer})
    st.rerun()
elif user_input:
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    st.session_state.chat_history.append({"role": "assistant", "content": "Please load precomputed embeddings/chunks and enter your Gemini API key to enable Q&A."})
    st.rerun() 